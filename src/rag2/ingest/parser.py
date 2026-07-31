"""
RAG² 多格式文档解析器。

按扩展名分派到专用解析器：
  PDF  → pypdfium2（文本）+ pdfplumber（表格）
  DOCX → mammoth（语义 markdown）
  HTML → trafilatura（正文提取）
  MD   → 原文读取（markdown-it-py 用于结构分析）
  PPTX → markitdown（如已安装）
  XLSX → openpyxl（表格 markdown）
  CSV  → pandas
  TXT  → 直接读

所有解析器输出统一的 ParsedDoc 格式。
"""
from __future__ import annotations

import hashlib
import logging
import mimetypes
import re
import time
from pathlib import Path
from typing import Optional

from rag2.ingest.models import ParsedDoc, file_hash

logger = logging.getLogger(__name__)

# 支持的扩展名 → MIME 类型映射
EXT_MIME = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".html": "text/html",
    ".htm": "text/html",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".csv": "text/csv",
    ".txt": "text/plain",
    ".rst": "text/x-rst",
}

SUPPORTED_EXTS = set(EXT_MIME.keys())


class DocumentParser:
    """多格式文档解析器。

    用法：
        parser = DocumentParser()
        doc = parser.parse(Path("report.pdf"))
        docs = parser.parse_dir(Path("./documents"))
    """

    def parse(self, file_path: Path) -> ParsedDoc:
        """解析单个文件，返回 ParsedDoc。"""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = file_path.suffix.lower()
        if ext not in EXT_MIME:
            raise ValueError(f"不支持的格式: {ext} (文件: {file_path})")

        mime = EXT_MIME[ext]
        doc_id = file_hash(file_path)
        t0 = time.time()

        # 分派到专用解析器
        if ext == ".pdf":
            text, sections, meta = self._parse_pdf(file_path)
        elif ext == ".docx":
            text, sections, meta = self._parse_docx(file_path)
        elif ext in (".html", ".htm"):
            text, sections, meta = self._parse_html(file_path)
        elif ext in (".md", ".markdown"):
            text, sections, meta = self._parse_markdown(file_path)
        elif ext == ".pptx":
            text, sections, meta = self._parse_pptx(file_path)
        elif ext == ".xlsx":
            text, sections, meta = self._parse_xlsx(file_path)
        elif ext == ".csv":
            text, sections, meta = self._parse_csv(file_path)
        elif ext == ".doc":
            text, sections, meta = self._parse_doc_legacy(file_path)
        else:  # .txt, .rst
            text, sections, meta = self._parse_text(file_path)

        parse_time = time.time() - t0
        title = self._extract_title(text, file_path, sections)
        meta["parser"] = f"rag2-ingest-{ext}"
        meta["parse_time_s"] = round(parse_time, 2)
        meta["n_chars"] = len(text)
        meta["n_sections"] = len(sections)

        logger.info("解析 %s (%s): %d 字符, %d 段落, %.1fs",
                     file_path.name, ext, len(text), len(sections), parse_time)

        return ParsedDoc(
            doc_id=doc_id,
            source_path=str(file_path),
            filename=file_path.name,
            mime_type=mime,
            title=title,
            text=text,
            sections=sections,
            metadata=meta,
        )

    def parse_dir(self, dir_path: Path, recursive: bool = True) -> list[ParsedDoc]:
        """批量解析目录。跳过不支持的格式和隐藏文件。"""
        dir_path = Path(dir_path)
        if not dir_path.is_dir():
            raise NotADirectoryError(f"不是目录: {dir_path}")

        glob_pattern = "**/*" if recursive else "*"
        docs = []
        skipped = 0
        for f in sorted(dir_path.glob(glob_pattern)):
            if not f.is_file():
                continue
            if f.name.startswith(".") or f.suffix.lower() not in SUPPORTED_EXTS:
                skipped += 1
                continue
            try:
                doc = self.parse(f)
                docs.append(doc)
            except Exception as e:
                logger.warning("解析失败 %s: %s", f, e)
                skipped += 1

        logger.info("目录 %s: 解析 %d 文件, 跳过 %d", dir_path.name, len(docs), skipped)
        return docs

    # ── 各格式解析器 ──────────────────────────────────

    def _parse_pdf(self, path: Path) -> tuple[str, list[dict], dict]:
        """PDF: pypdfium2 提取文本 + pdfplumber 提取表格。"""
        import pypdfium2 as pdfium

        text_parts = []
        sections = []
        all_tables = []

        pdf = pdfium.PdfDocument(str(path))
        n_pages = len(pdf)

        for page_idx in range(n_pages):
            page = pdf[page_idx]
            textpage = page.get_textpage()
            page_text = textpage.get_text_range()
            textpage.close()
            page.close()

            if page_text.strip():
                text_parts.append(page_text)
                # 按标题行分段（简单启发式：短行可能是标题）
                for line in page_text.split("\n"):
                    line = line.strip()
                    if line and len(line) < 80 and not line.endswith("."):
                        sections.append({"heading": line, "text": "", "page": page_idx + 1, "level": 2})
                    elif line:
                        if sections:
                            sections[-1]["text"] += line + " "

        pdf.close()

        # pdfplumber 提取表格（如果有的话）
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf_pc:
                for page_idx, page in enumerate(pdf_pc.pages[:20]):  # 限制前20页
                    tables = page.extract_tables()
                    for table in tables:
                        if table and len(table) > 1:
                            # 转 markdown 表格
                            header = table[0]
                            rows = table[1:]
                            md = self._table_to_markdown(header, rows)
                            all_tables.append({"page": page_idx + 1, "table": md})
        except Exception as e:
            logger.debug("pdfplumber 表格提取跳过 %s: %s", path.name, e)

        text = "\n\n".join(text_parts)
        if all_tables:
            tables_md = "\n\n".join(f"[Table p{t['page']}]\n{t['table']}" for t in all_tables)
            text += "\n\n" + tables_md

        meta = {"pages": n_pages, "tables": len(all_tables)}
        return text, sections, meta

    def _parse_docx(self, path: Path) -> tuple[str, list[dict], dict]:
        """DOCX: mammoth 转 markdown，保留标题层级。"""
        import mammoth

        with open(path, "rb") as f:
            result = mammoth.convert_to_markdown(f)
            text = result.value
            messages = result.messages

        # 从 markdown 提取 sections
        sections = self._markdown_to_sections(text)

        meta = {"warnings": len(messages)}
        return text, sections, meta

    def _parse_html(self, path: Path) -> tuple[str, list[dict], dict]:
        """HTML: trafilatura 提取正文。"""
        import trafilatura

        html_content = path.read_text(encoding="utf-8", errors="ignore")

        # 提取正文 + metadata
        extracted = trafilatura.extract(
            html_content, output_format="markdown",
            with_metadata=True, include_tables=True,
        )
        if extracted:
            text = extracted
        else:
            # fallback: 简单标签清理
            text = re.sub(r"<[^>]+>", " ", html_content)
            text = re.sub(r"\s+", " ", text).strip()

        # metadata
        meta_dict = trafilatura.extract(html_content, output_format="json", with_metadata=True)
        meta = {}
        if meta_dict:
            try:
                import json
                d = json.loads(meta_dict)
                meta = {
                    "author": d.get("author", ""),
                    "url": d.get("url", ""),
                    "date": d.get("date", ""),
                }
            except Exception:
                pass

        sections = self._markdown_to_sections(text)
        return text, sections, meta

    def _parse_markdown(self, path: Path) -> tuple[str, list[dict], dict]:
        """Markdown: 原文 + 结构化 sections。"""
        text = path.read_text(encoding="utf-8", errors="ignore")
        sections = self._markdown_to_sections(text)
        meta = {}
        return text, sections, meta

    def _parse_pptx(self, path: Path) -> tuple[str, list[dict], dict]:
        """PPTX: 尝试 markitdown，fallback 到 python-pptx。"""
        try:
            from markitdown import MarkItDown
            md = MarkItDown()
            result = md.convert(str(path))
            text = result.text_content
        except ImportError:
            # fallback: python-pptx
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"## Slide {i+1}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"*Notes: {notes}*")
                parts.append("")
            text = "\n\n".join(parts)

        sections = self._markdown_to_sections(text)
        meta = {}
        return text, sections, meta

    def _parse_xlsx(self, path: Path) -> tuple[str, list[dict], dict]:
        """XLSX: openpyxl 转 markdown 表格。"""
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if rows:
                header = [str(c) if c is not None else "" for c in rows[0]]
                data = [[str(c) if c is not None else "" for c in row] for row in rows[1:50]]  # 限50行
                md = self._table_to_markdown(header, data)
                parts.append(f"## {sheet_name}\n\n{md}")

        wb.close()
        text = "\n\n".join(parts)
        sections = [{"heading": s, "text": "", "page": None, "level": 2}
                    for s in wb.sheetnames]
        meta = {"sheets": len(wb.sheetnames)}
        return text, sections, meta

    def _parse_csv(self, path: Path) -> tuple[str, list[dict], dict]:
        """CSV: pandas 读表格转 markdown。"""
        import pandas as pd
        df = pd.read_csv(path, nrows=100)  # 限100行
        text = df.to_markdown(index=False)
        sections = []
        meta = {"rows": len(df), "cols": len(df.columns)}
        return text, sections, meta

    def _parse_doc_legacy(self, path: Path) -> tuple[str, list[dict], dict]:
        """旧版 .doc 格式（有限支持）。"""
        # 尝试 antiword 或直接二进制提取
        try:
            import subprocess
            result = subprocess.run(["antiword", str(path)], capture_output=True, text=True, timeout=10)
            if result.returncode == 0:
                text = result.stdout
            else:
                raise RuntimeError("antiword failed")
        except Exception:
            # fallback: 粗暴提取可读文本
            raw = path.read_bytes()
            text = raw.decode("utf-8", errors="ignore")
            text = re.sub(r"[^\x20-\x7e\n]", " ", text)
            text = re.sub(r"\s+", " ", text).strip()

        sections = []
        meta = {"note": "legacy .doc format, limited extraction"}
        return text, sections, meta

    def _parse_text(self, path: Path) -> tuple[str, list[dict], dict]:
        """TXT/RST: 直接读。"""
        text = path.read_text(encoding="utf-8", errors="ignore")
        sections = self._markdown_to_sections(text)
        meta = {}
        return text, sections, meta

    # ── 工具方法 ──────────────────────────────────────

    @staticmethod
    def _table_to_markdown(header: list[str], rows: list[list[str]]) -> str:
        """表格转 markdown 格式。"""
        if not header:
            return ""
        lines = ["| " + " | ".join(str(h) for h in header) + " |"]
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows:
            # 补齐列数
            while len(row) < len(header):
                row.append("")
            lines.append("| " + " | ".join(str(c) for c in row[:len(header)]) + " |")
        return "\n".join(lines)

    @staticmethod
    def _markdown_to_sections(text: str) -> list[dict]:
        """从 markdown 文本提取 sections（按标题层级切分）。"""
        sections = []
        current_heading = ""
        current_text = []
        current_level = 0

        for line in text.split("\n"):
            # 匹配 markdown 标题
            m = re.match(r"^(#{1,6})\s+(.+)$", line)
            if m:
                # 保存前一个 section
                if current_heading or current_text:
                    sections.append({
                        "heading": current_heading,
                        "text": "\n".join(current_text).strip(),
                        "page": None,
                        "level": current_level,
                    })
                current_heading = m.group(2).strip()
                current_level = len(m.group(1))
                current_text = []
            else:
                current_text.append(line)

        # 最后一个 section
        if current_heading or current_text:
            sections.append({
                "heading": current_heading,
                "text": "\n".join(current_text).strip(),
                "page": None,
                "level": current_level,
            })

        # 过滤空 sections
        return [s for s in sections if s["text"] or s["heading"]]

    @staticmethod
    def _extract_title(text: str, path: Path, sections: list[dict]) -> str:
        """提取文档标题：第一个标题 > 文件名。"""
        # 尝试从 sections 找第一个 level-1 标题
        for s in sections:
            if s.get("level", 0) <= 1 and s["heading"]:
                return s["heading"]
        # fallback: 文件名（去扩展名）
        return path.stem.replace("_", " ").replace("-", " ").title()
